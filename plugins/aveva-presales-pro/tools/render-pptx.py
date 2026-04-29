#!/usr/bin/env python3
"""
Native PPTX rendering for aveva-presales-pro plugin.
Parses gamma-optimized markdown and produces a PowerPoint deck using python-pptx.

Slide-break convention: '---' on its own line.
Slide title: '## ' or '### ' heading at start of each section.
Bullets: '- ' lines.

Usage:
    python3 tools/render-pptx.py <input.md> <output.pptx> [--config distributor.config.yaml]

Requires: python-pptx, pyyaml
Install: pip install python-pptx pyyaml
"""
import sys
import os
import argparse
import yaml
import re
from datetime import datetime
import json

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print("ERROR: Missing dependency. Run: pip install python-pptx pyyaml", file=sys.stderr)
    print(f"  Detail: {e}", file=sys.stderr)
    sys.exit(2)


def read_plugin_version(root="."):
    """Read version from .claude-plugin/plugin.json. Falls back to 'unknown'."""
    candidates = [
        os.path.join(root, ".claude-plugin", "plugin.json"),
        os.path.join(os.path.dirname(__file__), "..", ".claude-plugin", "plugin.json"),
    ]
    for p in candidates:
        try:
            return json.load(open(p)).get("version", "unknown")
        except (OSError, json.JSONDecodeError):
            continue
    return "unknown"


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_config(config_path):
    if not os.path.exists(config_path):
        return {"brand": {"primary_color": "#0F4C81", "accent_color": "#00B5B0"},
                "distributor": {"name": "Distributor"}}
    with open(config_path) as f:
        return yaml.safe_load(f)


def parse_gamma_markdown(text):
    """Returns list of slides; each slide is dict with title and bullets."""
    # Split on lines that are exactly '---'
    raw_slides = re.split(r"^---\s*$", text, flags=re.MULTILINE)
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        slide = {"title": "", "subtitle": "", "bullets": [], "body_paragraphs": []}
        lines = raw.split("\n")
        in_bullets = False
        body_buf = []
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                if body_buf:
                    slide["body_paragraphs"].append(" ".join(body_buf).strip())
                    body_buf = []
                continue
            # H1 = deck title (first slide); but treat as title if encountered
            m1 = re.match(r"^#\s+(.+)$", line_stripped)
            m2 = re.match(r"^##\s+(.+)$", line_stripped)
            m3 = re.match(r"^###\s+(.+)$", line_stripped)
            if m1 or m2 or m3:
                if body_buf:
                    slide["body_paragraphs"].append(" ".join(body_buf).strip())
                    body_buf = []
                title = (m1 or m2 or m3).group(1).strip()
                if not slide["title"]:
                    slide["title"] = title
                else:
                    slide["subtitle"] = title
                continue
            mb = re.match(r"^[-*]\s+(.+)$", line_stripped)
            if mb:
                if body_buf:
                    slide["body_paragraphs"].append(" ".join(body_buf).strip())
                    body_buf = []
                slide["bullets"].append(mb.group(1).strip())
                continue
            # ordinary text
            body_buf.append(line_stripped)
        if body_buf:
            slide["body_paragraphs"].append(" ".join(body_buf).strip())
        slides.append(slide)
    return slides


def build_deck(slides, output_path, config, plugin_version="unknown"):
    brand = config.get("brand", {})
    primary = brand.get("primary_color", "#0F4C81")
    accent = brand.get("accent_color", "#00B5B0")
    distributor = config.get("distributor", {}).get("name", "Distributor")

    template_path = brand.get("pptx_template")
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    primary_rgb = hex_to_rgb(primary)
    accent_rgb = hex_to_rgb(accent)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    title_layout = prs.slide_layouts[0]

    for idx, slide_data in enumerate(slides):
        if idx == 0:
            # Title slide
            slide = prs.slides.add_slide(title_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"] or "Untitled"
                for p in slide.shapes.title.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = primary_rgb
                        run.font.bold = True
            try:
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = f"{distributor} | {datetime.now().strftime('%Y-%m-%d')}"
            except Exception:
                pass
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Accent bar at top
        accent_box = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
        accent_box.fill.solid()
        accent_box.fill.fore_color.rgb = accent_rgb
        accent_box.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), prs.slide_width - Inches(1), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = slide_data["title"] or "Section"
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = primary_rgb

        # Body — bullets and paragraphs
        body_top = Inches(1.4)
        body_box = slide.shapes.add_textbox(Inches(0.5), body_top, prs.slide_width - Inches(1), prs.slide_height - body_top - Inches(0.5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True

        first = True
        for bullet in slide_data["bullets"]:
            p = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
            first = False
            p.text = "• " + bullet
            p.font.size = Pt(16)
        for para in slide_data["body_paragraphs"]:
            if para and para != slide_data["title"]:
                p = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
                first = False
                p.text = para
                p.font.size = Pt(14)

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.3), prs.slide_width - Inches(1), Inches(0.25))
        ftf = footer.text_frame
        ftf.text = f"{distributor} | aveva-presales-pro v{plugin_version}"
        for p in ftf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor(150, 150, 150)

    prs.save(output_path)
    print(f"Rendered: {output_path} ({len(slides)} slides) [plugin v{plugin_version}]")


def render(input_path, output_path, config_path):
    config = load_config(config_path)
    with open(input_path) as f:
        text = f.read()
    slides = parse_gamma_markdown(text)
    if not slides:
        print("ERROR: No slides parsed from input. Check that input has --- separators.", file=sys.stderr)
        sys.exit(1)
    build_deck(slides, output_path, config, read_plugin_version(os.path.dirname(os.path.dirname(os.path.abspath(input_path)))))


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Render gamma-markdown to PPTX")
    parser.add_argument("input", help="Input Markdown file")
    parser.add_argument("output", help="Output PPTX file")
    parser.add_argument("--config", default="distributor.config.yaml", help="Distributor config")
    args = parser.parse_args()
    render(args.input, args.output, args.config)
